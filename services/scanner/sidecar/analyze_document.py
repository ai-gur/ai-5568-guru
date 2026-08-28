#!/usr/bin/env python3
"""
IS 5568 Part 2 document analyser.

Reads a JSON job list on stdin, writes a JSON facts list on stdout. One process
handles the whole batch so the Python interpreter and the heavy PDF libraries
are started once per scan rather than once per file.

    echo '[{"path": "a.pdf", "kind": "pdf", "url": "..."}]' | python analyze_document.py

This extracts *facts*, never verdicts. Deciding whether a document passes is the
Node verdict layer's job, so that web pages and documents go through the same
precedence rules and the same NA discipline.

Part 2 defines its own large-text thresholds (section 3.6) which differ from the
web ones: in word-processing documents, "large" is 14pt bold or 18pt regular,
not the 18.5px/24px used for web pages.
"""

from __future__ import annotations

import json
import os
import re
import sys
import traceback
import zipfile
from typing import Any

# Section 3.6 of IS 5568 part 2.
LARGE_PT = 14.0          # large when bold
LARGE_PT_REGULAR = 18.0  # large when not bold
MIN_CONTRAST_NORMAL = 4.5
MIN_CONTRAST_LARGE = 3.0

HEBREW_RE = re.compile(r"[֐-׿]")
EMAIL_RE = re.compile(r"[\w.+-]+@[\w-]+\.[\w.]+")
URL_TEXT_RE = re.compile(r"^https?://", re.I)

GENERIC_LINK_HE = {
    "לחץ כאן", "לחצו כאן", "כאן", "קרא עוד", "קראו עוד", "למידע נוסף",
    "לפרטים", "עוד", "המשך", "click here", "read more", "more", "here", "link",
}

# Phrases that point at a sense rather than naming the target. No \b around the
# Hebrew alternatives: Python's \b is also ASCII-oriented for these scripts.
SENSORY_RE = re.compile(
    r"(?:בטבלה\s+(?:מימין|משמאל|למעלה|למטה)"
    r"|(?:הכפתור|הקישור|התיבה|הריבוע|העיגול|השדה)\s+(?:האדום|הירוק|הכחול|הצהוב|הכתום)"
    r"|המסומן\s+ב(?:אדום|ירוק|כחול|צהוב)"
    r"|(?:ראה|ראו|עיין|עיינו)\s+(?:בתרשים|בטבלה|באיור)\s+(?:מימין|משמאל|למעלה|למטה)"
    r"|בעיגול\s+האדום"
    r")"
)

COMPLEX_HINT_RE = re.compile(
    r"(?:תרשים|גרף|דיאגרמה|נוסחה|משוואה|אינפוגרפיק|chart|graph|diagram|formula|equation)",
    re.I,
)


def blank_facts(kind: str, url: str, path: str) -> dict[str, Any]:
    """Every key the Node side may read, so no consumer has to guard."""
    return {
        "kind": kind,
        "url": url,
        "fileName": os.path.basename(path),
        "bytes": os.path.getsize(path) if os.path.exists(path) else 0,
        "title": None,
        "language": None,
        "tagged": None,
        "displayDocTitle": None,
        "pageCount": 0,
        "images": [],
        "headings": [],
        "lists": [],
        "tables": [],
        "links": [],
        "textRuns": [],
        "contrastFailures": [],
        "readingOrderIssues": [],
        "sensoryPhrases": [],
        "complexInfo": [],
        "scannedPages": [],
        "textImages": [],
        "colouredRuns": [],
        "textLength": 0,
        "counts": {},
        "notes": [],
        "error": None,
    }


def finalize(f: dict[str, Any]) -> dict[str, Any]:
    f["counts"] = {
        "images": len(f["images"]),
        "headings": len(f["headings"]),
        "lists": len(f["lists"]),
        "tables": len(f["tables"]),
        "links": len(f["links"]),
        "paragraphs": len(f["textRuns"]),
        "textLength": f["textLength"],
        "sensoryPhrases": len(f["sensoryPhrases"]),
        "complexInfo": len(f["complexInfo"]),
        "scannedPages": len(f["scannedPages"]),
        "textImages": len(f["textImages"]),
        "colouredRuns": len(f["colouredRuns"]),
        "contrastFailures": len(f["contrastFailures"]),
    }
    return f


def is_large_text(size_pt: float, bold: bool) -> bool:
    return size_pt >= (LARGE_PT if bold else LARGE_PT_REGULAR)


def relative_luminance(rgb: tuple[float, float, float]) -> float:
    def channel(v: float) -> float:
        v = v / 255.0
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4

    r, g, b = (channel(c) for c in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(fg: tuple[float, float, float], bg: tuple[float, float, float]) -> float:
    l1, l2 = relative_luminance(fg), relative_luminance(bg)
    hi, lo = max(l1, l2), min(l1, l2)
    return round((hi + 0.05) / (lo + 0.05), 2)


def int_to_rgb(value: int) -> tuple[int, int, int]:
    return ((value >> 16) & 255, (value >> 8) & 255, value & 255)


# ── PDF ──────────────────────────────────────────────────────────────────────

def analyse_pdf(path: str, url: str) -> dict[str, Any]:
    f = blank_facts("pdf", url, path)

    try:
        import pikepdf
    except ImportError:
        pikepdf = None  # type: ignore[assignment]

    if pikepdf is not None:
        try:
            with pikepdf.open(path) as pdf:
                root = pdf.Root
                mark_info = root.get("/MarkInfo")
                f["tagged"] = bool(mark_info and mark_info.get("/Marked"))
                f["language"] = str(root.get("/Lang")) if root.get("/Lang") else None

                viewer_prefs = root.get("/ViewerPreferences")
                f["displayDocTitle"] = (
                    bool(viewer_prefs.get("/DisplayDocTitle")) if viewer_prefs and "/DisplayDocTitle" in viewer_prefs else False
                )

                with pdf.open_metadata() as meta:
                    f["title"] = meta.get("dc:title") or None
                if not f["title"]:
                    info = pdf.docinfo
                    title = info.get("/Title") if info else None
                    f["title"] = str(title) if title else None

                # Walk the structure tree for heading/list/table/figure tags.
                struct_root = root.get("/StructTreeRoot")
                if struct_root is not None:
                    walk_pdf_struct(struct_root, f)
                elif f["tagged"]:
                    f["notes"].append("המסמך מסומן כ-Marked אך אין בו עץ תגיות (StructTreeRoot)")
        except Exception as exc:  # noqa: BLE001 — a broken PDF is a finding, not a crash
            f["notes"].append(f"pikepdf לא הצליח לקרוא את המבנה: {exc}")

    # Text, images and contrast come from PyMuPDF, which exposes span colours.
    try:
        import fitz  # PyMuPDF
    except ImportError:
        f["error"] = "pymupdf is not installed — run: python -m pip install -r requirements.txt"
        return finalize(f)

    try:
        doc = fitz.open(path)
    except Exception as exc:  # noqa: BLE001
        f["error"] = f"לא ניתן לפתוח את קובץ ה-PDF: {exc}"
        return finalize(f)

    f["pageCount"] = doc.page_count
    if f["title"] is None:
        f["title"] = (doc.metadata or {}).get("title") or None

    total_text = 0
    for page_index in range(doc.page_count):
        page = doc[page_index]
        page_text = page.get_text("text") or ""
        total_text += len(page_text.strip())

        page_images = page.get_images(full=True)
        page_rect = page.rect
        page_area = max(page_rect.width * page_rect.height, 1)

        # A page with an image covering most of it and almost no text is a scan.
        if page_images and len(page_text.strip()) < 40:
            biggest = 0.0
            for img in page_images:
                try:
                    for rect in page.get_image_rects(img[0]):
                        biggest = max(biggest, rect.width * rect.height)
                except Exception:  # noqa: BLE001
                    continue
            if biggest / page_area > 0.5:
                f["scannedPages"].append({"page": page_index + 1, "coverage": round(biggest / page_area, 2)})

        for phrase in SENSORY_RE.finditer(page_text):
            f["sensoryPhrases"].append({"page": page_index + 1, "match": phrase.group(0)})

        for hint in COMPLEX_HINT_RE.finditer(page_text):
            if len(f["complexInfo"]) < 40:
                start = max(0, hint.start() - 60)
                f["complexInfo"].append(
                    {"page": page_index + 1, "match": hint.group(0), "context": page_text[start : hint.end() + 60].replace("\n", " ")}
                )

        # Text spans, with size/bold/colour, for the contrast check.
        try:
            raw = page.get_text("dict")
        except Exception:  # noqa: BLE001
            raw = {"blocks": []}

        page_bg = (255, 255, 255)  # PDF pages are white unless painted otherwise
        for block in raw.get("blocks", []):
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = (span.get("text") or "").strip()
                    if not text:
                        continue
                    size = float(span.get("size", 0))
                    flags = int(span.get("flags", 0))
                    bold = bool(flags & 2 ** 4)
                    colour = int_to_rgb(int(span.get("color", 0)))
                    large = is_large_text(size, bold)
                    if len(f["textRuns"]) < 400:
                        f["textRuns"].append(
                            {
                                "page": page_index + 1,
                                "text": text[:120],
                                "sizePt": round(size, 1),
                                "bold": bold,
                                "large": large,
                                "color": f"rgb({colour[0]}, {colour[1]}, {colour[2]})",
                            }
                        )
                    if colour != (0, 0, 0) and len(f["colouredRuns"]) < 60:
                        f["colouredRuns"].append({"page": page_index + 1, "text": text[:80], "color": f"rgb{colour}"})

                    ratio = contrast_ratio(colour, page_bg)
                    required = MIN_CONTRAST_LARGE if large else MIN_CONTRAST_NORMAL
                    if ratio < required and len(f["contrastFailures"]) < 40:
                        f["contrastFailures"].append(
                            {
                                "page": page_index + 1,
                                "text": text[:80],
                                "ratio": ratio,
                                "required": required,
                                "sizePt": round(size, 1),
                                "bold": bold,
                                "large": large,
                            }
                        )

        for link in page.get_links():
            uri = link.get("uri")
            if not uri:
                continue
            # Link text is whatever sits inside the link rectangle.
            try:
                text = page.get_textbox(link["from"]).strip()
            except Exception:  # noqa: BLE001
                text = ""
            if len(f["links"]) < 200:
                f["links"].append(
                    {
                        "page": page_index + 1,
                        "uri": uri[:200],
                        "text": text[:120],
                        "generic": text.strip().lower() in GENERIC_LINK_HE,
                        "isUrlText": bool(URL_TEXT_RE.match(text.strip())),
                    }
                )

        for img in page_images:
            xref = img[0]
            try:
                rects = page.get_image_rects(xref)
                rect = rects[0] if rects else None
            except Exception:  # noqa: BLE001
                rect = None
            width, height = (round(rect.width), round(rect.height)) if rect else (0, 0)
            entry = {
                "page": page_index + 1,
                "xref": xref,
                "width": width,
                "height": height,
                # /Alt is read from the structure tree; matched up below.
                "alt": None,
                "decorative": width <= 24 and height <= 24,
            }
            if len(f["images"]) < 200:
                f["images"].append(entry)
            # A wide, short image is the classic "banner of text" shape.
            if width > 300 and height > 40 and width / max(height, 1) > 2.0 and len(f["textImages"]) < 40:
                f["textImages"].append({"page": page_index + 1, "width": width, "height": height})

    f["textLength"] = total_text
    doc.close()

    annotate_pdf_headings(f)

    # A PDF with pages but no extractable text is a scan even if no single page
    # tripped the coverage test above.
    if f["pageCount"] > 0 and total_text < 40 and not f["scannedPages"]:
        f["scannedPages"] = [{"page": i + 1, "coverage": 1.0} for i in range(min(f["pageCount"], 20))]
        f["notes"].append("לא נמצא טקסט הניתן לחילוץ במסמך כלל — סימן מובהק לקובץ סרוק")

    return finalize(f)


def annotate_pdf_headings(f: dict[str, Any]) -> None:
    """
    Supplies heading text, and finds headings that were never tagged.

    A PDF's structure tree records that something *is* an H2, but the words
    themselves live in the page content stream, so `/K` walking yields empty
    strings on most real files. Judgement rows ("does this heading describe its
    section?") need the actual text, and the highest-yield Part 2 finding —
    text styled as a heading but not tagged as one — needs the same font-size
    view. Both come from the spans already collected.
    """
    runs = f["textRuns"]
    if not runs:
        return

    sizes = [r["sizePt"] for r in runs if r.get("sizePt")]
    if not sizes:
        return
    # Body size = the most common span size on the document.
    modal = max(set(sizes), key=sizes.count)

    large_runs = [
        r
        for r in runs
        if r.get("sizePt")
        and len((r.get("text") or "").strip()) > 2
        and len((r.get("text") or "").strip()) < 120
        and (r["sizePt"] >= modal * 1.15 or (r.get("bold") and r["sizePt"] >= modal))
    ]

    tagged_headings = [h for h in f["headings"] if not h.get("fake")]

    # Fill in text for tagged headings, in document order, from the large spans.
    if tagged_headings and all(not (h.get("text") or "").strip() for h in tagged_headings):
        for heading, run in zip(tagged_headings, large_runs):
            heading["text"] = (run.get("text") or "")[:120]
            heading["textSource"] = "font-size-heuristic"
        f["notes"].append(
            "טקסט הכותרות ב-PDF שוחזר לפי גודל גופן: עץ התגיות מכיל את רמות הכותרות אך לא את מילותיהן."
        )

    # Large text with no heading tag at all — the document has visual headings
    # that assistive technology cannot navigate.
    if not tagged_headings and large_runs:
        for run in large_runs[:40]:
            f["headings"].append(
                {
                    "level": None,
                    "text": (run.get("text") or "")[:120],
                    "page": run.get("page"),
                    "fake": True,
                    "sizePt": run.get("sizePt"),
                    "bold": run.get("bold"),
                }
            )


def walk_pdf_struct(node: Any, f: dict[str, Any], depth: int = 0) -> None:
    """Collects heading/list/table/figure tags out of the PDF structure tree."""
    if depth > 40:
        return
    try:
        import pikepdf
    except ImportError:
        return

    kids = node.get("/K") if hasattr(node, "get") else None
    if kids is None:
        return
    if not isinstance(kids, (list, pikepdf.Array)):
        kids = [kids]

    for kid in kids:
        if not hasattr(kid, "get"):
            continue
        struct_type = kid.get("/S")
        name = str(struct_type).lstrip("/") if struct_type is not None else ""

        if re.fullmatch(r"H[1-6]", name):
            f["headings"].append({"level": int(name[1]), "text": extract_struct_text(kid)[:120]})
        elif name == "H":
            f["headings"].append({"level": None, "text": extract_struct_text(kid)[:120]})
        elif name == "L":
            f["lists"].append({"items": count_struct_children(kid, "/LI")})
        elif name == "Table":
            f["tables"].append(
                {
                    "headerCells": count_struct_children(kid, "/TH", recursive=True),
                    "dataCells": count_struct_children(kid, "/TD", recursive=True),
                }
            )
        elif name == "Figure":
            alt = kid.get("/Alt")
            f["images"].append(
                {
                    "page": None,
                    "xref": None,
                    "width": 0,
                    "height": 0,
                    "alt": str(alt) if alt is not None else None,
                    "decorative": False,
                    "fromStructTree": True,
                }
            )

        walk_pdf_struct(kid, f, depth + 1)


def extract_struct_text(node: Any, depth: int = 0) -> str:
    if depth > 8:
        return ""
    try:
        import pikepdf
    except ImportError:
        return ""
    out: list[str] = []
    kids = node.get("/K") if hasattr(node, "get") else None
    if kids is None:
        return ""
    if not isinstance(kids, (list, pikepdf.Array)):
        kids = [kids]
    for kid in kids:
        if isinstance(kid, pikepdf.String):
            out.append(str(kid))
        elif hasattr(kid, "get"):
            out.append(extract_struct_text(kid, depth + 1))
    return " ".join(x for x in out if x).strip()


def count_struct_children(node: Any, tag: str, recursive: bool = False, depth: int = 0) -> int:
    if depth > 20:
        return 0
    try:
        import pikepdf
    except ImportError:
        return 0
    kids = node.get("/K") if hasattr(node, "get") else None
    if kids is None:
        return 0
    if not isinstance(kids, (list, pikepdf.Array)):
        kids = [kids]
    total = 0
    for kid in kids:
        if not hasattr(kid, "get"):
            continue
        struct_type = kid.get("/S")
        if struct_type is not None and str(struct_type) == tag:
            total += 1
        if recursive:
            total += count_struct_children(kid, tag, True, depth + 1)
    return total


# ── DOCX ─────────────────────────────────────────────────────────────────────

def analyse_docx(path: str, url: str) -> dict[str, Any]:
    f = blank_facts("docx", url, path)
    try:
        import docx  # python-docx
    except ImportError:
        f["error"] = "python-docx is not installed — run: python -m pip install -r requirements.txt"
        return finalize(f)

    try:
        document = docx.Document(path)
    except Exception as exc:  # noqa: BLE001
        f["error"] = f"לא ניתן לפתוח את המסמך: {exc}"
        return finalize(f)

    core = document.core_properties
    f["title"] = (core.title or "").strip() or None
    f["language"] = (core.language or "").strip() or None
    f["tagged"] = True  # Office documents always carry style semantics

    text_total = 0
    manual_bullet = re.compile(r"^\s*(?:[-–—•*▪◦●·]|\d+[.)])\s+")

    for para in document.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        text_total += len(text)
        style = (para.style.name or "") if para.style else ""
        is_heading_style = style.startswith("Heading") or style.startswith("כותרת")
        level = None
        match = re.search(r"(\d+)", style)
        if is_heading_style and match:
            level = int(match.group(1))

        if is_heading_style:
            f["headings"].append({"level": level, "text": text[:120], "style": style})
        else:
            runs = [r for r in para.runs if (r.text or "").strip()]
            first = runs[0] if runs else None
            size_pt = float(first.font.size.pt) if first is not None and first.font.size else 0.0
            bold = bool(first.font.bold) if first is not None else False
            colour = None
            try:
                if first is not None and first.font.color is not None and first.font.color.rgb is not None:
                    colour = str(first.font.color.rgb)
            except Exception:  # noqa: BLE001
                colour = None

            if len(f["textRuns"]) < 400:
                f["textRuns"].append(
                    {
                        "text": text[:120],
                        "style": style,
                        "sizePt": size_pt,
                        "bold": bold,
                        "large": is_large_text(size_pt, bold) if size_pt else False,
                        "color": colour,
                    }
                )

            # A short, bold or enlarged paragraph that introduces the next one is
            # a heading that was never marked up — the commonest Part 2 failure.
            if text and len(text) < 120 and (bold or (size_pt and size_pt >= 14)):
                f["headings"].append(
                    {"level": None, "text": text[:120], "style": style, "fake": True, "sizePt": size_pt, "bold": bold}
                )

            if manual_bullet.match(text):
                f["lists"].append({"items": 1, "manual": True, "text": text[:80]})

            if colour and colour not in ("000000", "auto") and len(f["colouredRuns"]) < 60:
                f["colouredRuns"].append({"text": text[:80], "color": colour})

        for phrase in SENSORY_RE.finditer(text):
            f["sensoryPhrases"].append({"match": phrase.group(0), "context": text[:160]})
        for hint in COMPLEX_HINT_RE.finditer(text):
            if len(f["complexInfo"]) < 40:
                f["complexInfo"].append({"match": hint.group(0), "context": text[:160]})

    for numbered in document.paragraphs:
        style = (numbered.style.name or "") if numbered.style else ""
        if style in ("List Paragraph", "List Bullet", "List Number") or style.startswith("רשימה"):
            f["lists"].append({"items": 1, "manual": False, "style": style})

    for table in document.tables:
        rows = len(table.rows)
        cols = len(table.columns) if rows else 0
        # python-docx has no header-row API; the repeat-header flag in the XML is
        # the only reliable signal that a row was marked as a header.
        header_marked = False
        try:
            first_row = table.rows[0]._tr if rows else None  # noqa: SLF001
            if first_row is not None:
                header_marked = bool(first_row.findall(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblHeader"))
        except Exception:  # noqa: BLE001
            header_marked = False
        f["tables"].append(
            {
                "rows": rows,
                "cols": cols,
                "headerRowMarked": header_marked,
                "firstRowText": " | ".join(c.text.strip()[:20] for c in table.rows[0].cells) if rows else "",
            }
        )

    # Images and their alt text live in the drawing XML.
    f["images"].extend(docx_images(path))

    for rel in document.part.rels.values():
        if rel.reltype.endswith("/hyperlink"):
            f["links"].append({"uri": str(rel.target_ref)[:200], "text": "", "generic": False, "isUrlText": False})

    f["textLength"] = text_total
    return finalize(f)


def docx_images(path: str) -> list[dict[str, Any]]:
    """Alt text lives in <wp:docPr descr="…"> / title, which python-docx hides."""
    out: list[dict[str, Any]] = []
    try:
        with zipfile.ZipFile(path) as zf:
            for name in zf.namelist():
                if not (name.startswith("word/") and name.endswith(".xml")):
                    continue
                xml = zf.read(name).decode("utf-8", "ignore")
                for match in re.finditer(r"<wp:docPr\b([^>]*)/?>", xml):
                    attrs = match.group(1)
                    descr = re.search(r'descr="([^"]*)"', attrs)
                    title = re.search(r'title="([^"]*)"', attrs)
                    img_name = re.search(r'name="([^"]*)"', attrs)
                    decorative = "decorative" in attrs.lower()
                    out.append(
                        {
                            "name": img_name.group(1) if img_name else "",
                            "alt": (descr.group(1) if descr else None) or (title.group(1) if title else None),
                            "decorative": decorative,
                            "width": 0,
                            "height": 0,
                        }
                    )
                if len(out) > 200:
                    break
    except Exception:  # noqa: BLE001
        pass
    return out


# ── PPTX ─────────────────────────────────────────────────────────────────────

def analyse_pptx(path: str, url: str) -> dict[str, Any]:
    f = blank_facts("pptx", url, path)
    try:
        from pptx import Presentation
        from pptx.util import Pt  # noqa: F401
    except ImportError:
        f["error"] = "python-pptx is not installed — run: python -m pip install -r requirements.txt"
        return finalize(f)

    try:
        prs = Presentation(path)
    except Exception as exc:  # noqa: BLE001
        f["error"] = f"לא ניתן לפתוח את המצגת: {exc}"
        return finalize(f)

    core = prs.core_properties
    f["title"] = (core.title or "").strip() or None
    f["language"] = (core.language or "").strip() or None
    f["tagged"] = True
    f["pageCount"] = len(prs.slides)

    seen_titles: list[str] = []
    text_total = 0

    for index, slide in enumerate(prs.slides, start=1):
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:  # noqa: BLE001
            title_shape = None

        title_text = (title_shape.text or "").strip() if title_shape is not None else ""
        f["headings"].append(
            {
                "level": 1,
                "text": title_text[:120],
                "slide": index,
                "missing": not title_text,
                "duplicate": bool(title_text) and title_text in seen_titles,
            }
        )
        if title_text:
            seen_titles.append(title_text)

        # Reading order is the shape order in the tree; a title that is not
        # first is announced after body content.
        shape_order: list[str] = []
        for shape in slide.shapes:
            shape_order.append(shape.shape_type.__str__() if shape.shape_type is not None else "shape")

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(r.text or "" for r in para.runs).strip()
                    if not text:
                        continue
                    text_total += len(text)
                    first = para.runs[0] if para.runs else None
                    size_pt = float(first.font.size.pt) if first is not None and first.font.size else 0.0
                    bold = bool(first.font.bold) if first is not None else False
                    if len(f["textRuns"]) < 400:
                        f["textRuns"].append(
                            {
                                "slide": index,
                                "text": text[:120],
                                "sizePt": size_pt,
                                "bold": bold,
                                "large": is_large_text(size_pt, bold) if size_pt else False,
                            }
                        )
                    for phrase in SENSORY_RE.finditer(text):
                        f["sensoryPhrases"].append({"slide": index, "match": phrase.group(0)})
                    for hint in COMPLEX_HINT_RE.finditer(text):
                        if len(f["complexInfo"]) < 40:
                            f["complexInfo"].append({"slide": index, "match": hint.group(0), "context": text[:160]})

            if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                alt = ""
                try:
                    alt = shape._element._nvXxPr.cNvPr.get("descr", "") or ""  # noqa: SLF001
                except Exception:  # noqa: BLE001
                    alt = ""
                f["images"].append({"slide": index, "alt": alt or None, "decorative": False, "width": 0, "height": 0})

            if getattr(shape, "has_table", False):
                table = shape.table
                f["tables"].append(
                    {
                        "slide": index,
                        "rows": len(table.rows),
                        "cols": len(table.columns),
                        "headerRowMarked": bool(getattr(table, "first_row", False)),
                    }
                )

            if getattr(shape, "has_chart", False):
                f["complexInfo"].append({"slide": index, "match": "chart", "context": "תרשים במצגת"})

        if title_shape is not None and slide.shapes and slide.shapes[0] != title_shape:
            f["readingOrderIssues"].append(
                {"slide": index, "issue": "כותרת השקופית אינה הראשונה בסדר הקריאה של השקופית"}
            )

    f["textLength"] = text_total
    return finalize(f)


# ── XLSX ─────────────────────────────────────────────────────────────────────

def analyse_xlsx(path: str, url: str) -> dict[str, Any]:
    f = blank_facts("xlsx", url, path)
    try:
        import openpyxl
    except ImportError:
        f["error"] = "openpyxl is not installed — run: python -m pip install -r requirements.txt"
        return finalize(f)

    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception as exc:  # noqa: BLE001
        f["error"] = f"לא ניתן לפתוח את הגיליון: {exc}"
        return finalize(f)

    props = wb.properties
    f["title"] = (props.title or "").strip() or None
    f["language"] = (props.language or "").strip() or None
    f["tagged"] = True
    f["pageCount"] = len(wb.sheetnames)

    text_total = 0
    default_name = re.compile(r"^(Sheet|גיליון)\s*\d*$", re.I)

    for sheet in wb.worksheets:
        f["headings"].append(
            {
                "level": 1,
                "text": sheet.title,
                "sheet": sheet.title,
                # A default sheet name carries no information about its contents.
                "generic": bool(default_name.match(sheet.title.strip())),
            }
        )

        merged = list(sheet.merged_cells.ranges)
        rows = sheet.max_row or 0
        cols = sheet.max_column or 0
        first_row_values: list[str] = []
        if rows:
            for cell in next(sheet.iter_rows(min_row=1, max_row=1), []):
                if cell.value is not None:
                    first_row_values.append(str(cell.value)[:30])

        f["tables"].append(
            {
                "sheet": sheet.title,
                "rows": rows,
                "cols": cols,
                "mergedCells": len(merged),
                "firstRowText": " | ".join(first_row_values),
                # A defined table object is how Excel exposes a real header row.
                "headerRowMarked": bool(getattr(sheet, "tables", None)),
            }
        )

        for row in sheet.iter_rows(max_row=min(rows, 400)):
            for cell in row:
                if cell.value is None:
                    continue
                value = str(cell.value)
                text_total += len(value)
                if len(f["textRuns"]) < 400:
                    f["textRuns"].append({"sheet": sheet.title, "cell": cell.coordinate, "text": value[:80]})
                for phrase in SENSORY_RE.finditer(value):
                    f["sensoryPhrases"].append({"sheet": sheet.title, "cell": cell.coordinate, "match": phrase.group(0)})
                try:
                    fill = cell.fill
                    if fill is not None and fill.fgColor is not None and fill.fgColor.rgb not in (None, "00000000") and len(f["colouredRuns"]) < 60:
                        f["colouredRuns"].append({"sheet": sheet.title, "cell": cell.coordinate, "color": str(fill.fgColor.rgb)})
                except Exception:  # noqa: BLE001
                    pass

        for image in getattr(sheet, "_images", []):  # noqa: SLF001
            alt = None
            try:
                alt = getattr(image, "desc", None) or None
            except Exception:  # noqa: BLE001
                alt = None
            f["images"].append({"sheet": sheet.title, "alt": alt, "decorative": False, "width": 0, "height": 0})

        if getattr(sheet, "_charts", None):  # noqa: SLF001
            for _ in sheet._charts:  # noqa: SLF001
                f["complexInfo"].append({"sheet": sheet.title, "match": "chart", "context": "תרשים בגיליון"})

    f["textLength"] = text_total
    return finalize(f)


# ── plain text ───────────────────────────────────────────────────────────────

def analyse_txt(path: str, url: str) -> dict[str, Any]:
    f = blank_facts("txt", url, path)
    raw = b""
    try:
        with open(path, "rb") as handle:
            raw = handle.read(2_000_000)
    except Exception as exc:  # noqa: BLE001
        f["error"] = f"לא ניתן לקרוא את הקובץ: {exc}"
        return finalize(f)

    encoding = "utf-8"
    try:
        text = raw.decode("utf-8")
    except UnicodeDecodeError:
        # Legacy Hebrew files are usually cp1255; a mis-decoded file is itself
        # an accessibility problem because screen readers get mojibake.
        encoding = "cp1255"
        text = raw.decode("cp1255", "replace")
        f["notes"].append("הקובץ אינו בקידוד UTF-8. קידוד לא סטנדרטי עלול לגרום להצגה שגויה של עברית.")

    f["language"] = "he" if HEBREW_RE.search(text) else None
    f["tagged"] = False
    f["textLength"] = len(text.strip())
    f["notes"].append(f"קידוד שזוהה: {encoding}")

    lines = text.splitlines()
    for index, line in enumerate(lines[:500], start=1):
        stripped = line.strip()
        if not stripped:
            continue
        if len(f["textRuns"]) < 400:
            f["textRuns"].append({"line": index, "text": stripped[:120]})
        # Plain text has no semantics; the standard accepts a consistent,
        # continuous manual numbering as the hierarchy mechanism.
        if re.match(r"^\d+(\.\d+)*[.)]?\s+\S", stripped):
            f["headings"].append({"level": stripped.count(".") + 1, "text": stripped[:120], "line": index})
        elif re.match(r"^\s*[-–—•*]\s+\S", line):
            f["lists"].append({"items": 1, "manual": True, "line": index})
        for phrase in SENSORY_RE.finditer(stripped):
            f["sensoryPhrases"].append({"line": index, "match": phrase.group(0)})

    for match in EMAIL_RE.finditer(text):
        if len(f["links"]) < 100:
            f["links"].append({"uri": f"mailto:{match.group(0)}", "text": match.group(0), "generic": False, "isUrlText": False})

    return finalize(f)


ANALYSERS = {
    "pdf": analyse_pdf,
    "docx": analyse_docx,
    "pptx": analyse_pptx,
    "xlsx": analyse_xlsx,
    "txt": analyse_txt,
}


def main() -> int:
    try:
        jobs = json.load(sys.stdin)
    except Exception as exc:  # noqa: BLE001
        json.dump({"error": f"invalid job JSON on stdin: {exc}"}, sys.stdout)
        return 1

    results: list[dict[str, Any]] = []
    for job in jobs:
        kind = job.get("kind", "")
        path = job.get("path", "")
        url = job.get("url", "")
        analyser = ANALYSERS.get(kind)
        if analyser is None:
            facts = blank_facts(kind or "unknown", url, path)
            facts["error"] = f"unsupported document kind: {kind!r}"
            results.append(finalize(facts))
            continue
        try:
            results.append(analyser(path, url))
        except Exception:  # noqa: BLE001 — one bad file must not lose the batch
            facts = blank_facts(kind, url, path)
            facts["error"] = f"analysis crashed: {traceback.format_exc(limit=3)}"
            results.append(finalize(facts))

    json.dump(results, sys.stdout, ensure_ascii=False)
    return 0


if __name__ == "__main__":
    sys.exit(main())
